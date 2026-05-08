import io
import re
import textwrap
from pathlib import Path

import fitz  # pymupdf
import mammoth
import markdownify
from flask import Flask, jsonify, render_template, request
from pptx import Presentation
from pptx.util import Pt

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 50 * 1024 * 1024  # 50 MB

ALLOWED_EXTENSIONS = {
    ".docx", ".doc",
    ".pptx", ".ppt",
    ".pdf",
    ".txt", ".md",
}


def convert_docx(file_bytes: bytes) -> str:
    result = mammoth.convert_to_html(io.BytesIO(file_bytes))
    html = result.value
    md = markdownify.markdownify(html, heading_style="ATX", bullets="-")
    return _clean(md)


def convert_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    lines: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        title_text = ""
        body_parts: list[str] = []

        # Collect shapes sorted top-to-bottom so title comes first
        shapes = sorted(slide.shapes, key=lambda s: s.top if s.top is not None else 0)

        for shape in shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue

            # Heuristic: treat the first text shape (or placeholder named title) as heading
            is_title = (
                shape.name.lower().startswith("title")
                or (not title_text and shape == shapes[0])
            )

            if is_title and not title_text:
                title_text = text
            else:
                # Preserve bullet structure from paragraphs
                for para in shape.text_frame.paragraphs:
                    para_text = para.text.strip()
                    if not para_text:
                        continue
                    level = para.level or 0
                    indent = "  " * level
                    body_parts.append(f"{indent}- {para_text}")

        if title_text:
            lines.append(f"## Slide {i}: {title_text}")
        else:
            lines.append(f"## Slide {i}")

        if body_parts:
            lines.extend(body_parts)

        lines.append("")  # blank line between slides

    return "\n".join(lines)


def convert_pdf(file_bytes: bytes) -> str:
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    parts: list[str] = []

    for page_num, page in enumerate(doc, 1):
        blocks = page.get_text("dict")["blocks"]
        page_lines: list[str] = []

        for block in blocks:
            if block.get("type") != 0:  # 0 = text block
                continue
            for line in block.get("lines", []):
                spans = line.get("spans", [])
                if not spans:
                    continue

                text = "".join(s["text"] for s in spans).strip()
                if not text:
                    continue

                # Use font size to guess headings
                max_size = max(s["size"] for s in spans)
                flags = spans[0].get("flags", 0)
                bold = bool(flags & 2**4)

                if max_size >= 16 or (bold and max_size >= 13):
                    page_lines.append(f"## {text}")
                elif max_size >= 13 or bold:
                    page_lines.append(f"### {text}")
                else:
                    page_lines.append(text)

        if page_lines:
            parts.append(f"<!-- Page {page_num} -->")
            parts.extend(page_lines)
            parts.append("")

    doc.close()
    return "\n".join(parts)


def convert_text(file_bytes: bytes) -> str:
    return file_bytes.decode("utf-8", errors="replace")


def _clean(text: str) -> str:
    # Collapse 3+ blank lines into 2
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


CONVERTERS = {
    ".docx": convert_docx,
    ".doc":  convert_docx,
    ".pptx": convert_pptx,
    ".ppt":  convert_pptx,
    ".pdf":  convert_pdf,
    ".txt":  convert_text,
    ".md":   convert_text,
}


@app.route("/")
def index():
    return render_template("index.html")


@app.route("/convert", methods=["POST"])
def convert():
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400

    f = request.files["file"]
    if not f.filename:
        return jsonify({"error": "Empty filename"}), 400

    ext = Path(f.filename).suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        return jsonify({"error": f"Unsupported file type: {ext}"}), 415

    converter = CONVERTERS.get(ext)
    if not converter:
        return jsonify({"error": f"No converter for {ext}"}), 415

    try:
        data = f.read()
        markdown = converter(data)
        return jsonify({"markdown": markdown, "filename": f.filename})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


if __name__ == "__main__":
    app.run(debug=True, port=5001)
