"""FastAPI server for Refinery: distill a file, optionally redact, render to any format."""

from __future__ import annotations

import base64
import io
import json
import logging
from pathlib import Path
from typing import Any

from fastapi import FastAPI, File, Form, HTTPException, Request, UploadFile
from fastapi.responses import HTMLResponse, JSONResponse
from fastapi.staticfiles import StaticFiles

from refinery.engine import DEFAULT_ENTITIES, Redactor

log = logging.getLogger(__name__)

MAX_BYTES = 50 * 1024 * 1024  # 50 MB (Markify's limit; Redactor was 10 MB)
STATIC_DIR = Path(__file__).parent / "static"

ALLOWED_OUTPUT_FORMATS = {
    "original", "txt", "md", "pdf", "docx", "html", "rtf", "odt", "epub",
}

OUTPUT_MIME = {
    "txt": "text/plain; charset=utf-8",
    "md": "text/markdown; charset=utf-8",
    "html": "text/html; charset=utf-8",
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "rtf": "application/rtf",
    "odt": "application/vnd.oasis.opendocument.text",
    "epub": "application/epub+zip",
    "eml": "message/rfc822",
}

# Suffixes whose "original" preserves binary structure rather than re-rendering.
BINARY_ORIGINAL_SUFFIXES = {".docx", ".pptx", ".eml", ".msg", ".pdf"}

_redactor_singleton: Redactor | None = None


def get_redactor() -> Redactor:
    """Lazily build a single Redactor (loads spaCy on first use)."""
    global _redactor_singleton
    if _redactor_singleton is None:
        log.info("Initializing Redactor (loading spaCy model)...")
        _redactor_singleton = Redactor()
    return _redactor_singleton


def _build_redactor(
    mapping: dict[str, dict[str, str]] | None,
    entities: list[str] | None,
    threshold: float | None,
) -> Redactor:
    """Per-request Redactor that shares the loaded spaCy analyzer."""
    base = get_redactor()
    r = Redactor.__new__(Redactor)
    r._analyzer = base._analyzer
    r._faker_locale = base._faker_locale
    r.entities = list(entities) if entities else list(base.entities)
    r.score_threshold = threshold if threshold is not None else base.score_threshold
    r.mapping = json.loads(json.dumps(mapping)) if mapping else {}
    return r


def _redact_text(text: str, mapping, entities, threshold) -> tuple[str, dict]:
    r = _build_redactor(mapping, entities, threshold)
    return r.redact(text).text, r.mapping


def _parse_optional_json(raw: str | None, field_name: str) -> Any:
    if raw is None or raw == "":
        return None
    try:
        return json.loads(raw)
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=400, detail=f"{field_name} is not valid JSON: {e}")


def _parse_bool(value: Any, default: bool = False) -> bool:
    if value is None:
        return default
    if isinstance(value, bool):
        return value
    s = str(value).strip().lower()
    if s in ("1", "true", "yes", "on"):
        return True
    if s in ("0", "false", "no", "off", ""):
        return False
    return default


def create_app() -> FastAPI:
    app = FastAPI(title="Refinery", version="0.1.0")

    @app.exception_handler(Exception)
    async def _json_error_handler(request: Request, exc: Exception) -> JSONResponse:
        log.exception("Unhandled error processing %s %s", request.method, request.url.path)
        detail = f"{type(exc).__name__}: {exc}"
        return JSONResponse(status_code=500, content={"detail": detail})

    @app.get("/", response_class=HTMLResponse)
    def index() -> HTMLResponse:
        return HTMLResponse((STATIC_DIR / "index.html").read_text(encoding="utf-8"))

    @app.get("/api/entities")
    def list_entities() -> dict:
        return {"entities": DEFAULT_ENTITIES}

    @app.get("/api/capabilities")
    def capabilities() -> dict:
        from refinery import pandoc_render

        return {
            "pandoc": pandoc_render.is_pandoc_available(),
            "pdf_engine": pandoc_render.find_pdf_engine(),
        }

    @app.post("/api/process")
    async def process(
        file: UploadFile | None = File(default=None),
        text: str | None = Form(default=None),
        redact: str | None = Form(default="true"),
        load_mapping: str | None = Form(default=None),
        entities: str | None = Form(default=None),
        threshold: float | None = Form(default=None),
        output_format: str = Form(default="original"),
    ) -> dict:
        if file is None and not text:
            raise HTTPException(status_code=400, detail="Provide a file or text.")
        if output_format not in ALLOWED_OUTPUT_FORMATS:
            raise HTTPException(
                status_code=400,
                detail=f"output_format must be one of: {sorted(ALLOWED_OUTPUT_FORMATS)}",
            )

        do_redact = _parse_bool(redact, default=True)
        ent_list = _parse_optional_json(entities, "entities") if entities else None
        if ent_list is not None and not isinstance(ent_list, list):
            raise HTTPException(status_code=400, detail="entities must be a JSON array")
        loaded_map = _parse_optional_json(load_mapping, "load_mapping")
        if loaded_map is not None and not isinstance(loaded_map, dict):
            raise HTTPException(status_code=400, detail="load_mapping must be a JSON object")

        if file is not None and file.filename:
            data = await file.read()
            if len(data) > MAX_BYTES:
                raise HTTPException(
                    status_code=413,
                    detail=f"File exceeds {MAX_BYTES // (1024 * 1024)} MB.",
                )
            return _process_upload(
                file.filename, data, do_redact, loaded_map, ent_list, threshold, output_format,
            )

        return _process_text(text or "", do_redact, loaded_map, ent_list, threshold, output_format)

    @app.post("/api/reverse")
    def reverse(payload: dict) -> dict:
        from refinery.mapping import reverse_text

        body_text = payload.get("text", "")
        mapping = payload.get("mapping") or {}
        return {"text": reverse_text(body_text, mapping)}

    app.mount("/static", StaticFiles(directory=str(STATIC_DIR)), name="static")
    return app


def _process_text(
    text: str, do_redact: bool, loaded_map, entities, threshold, output_format: str,
) -> dict:
    """Pasted-text path: nothing to extract, just optionally redact and render."""
    if do_redact:
        body, mapping = _redact_text(text, loaded_map, entities, threshold)
    else:
        body, mapping = text, {}

    fmt = "txt" if output_format == "original" else output_format
    return _build_response(body, mapping, f"refinery.{fmt}", fmt)


def _process_upload(
    filename: str,
    data: bytes,
    do_redact: bool,
    loaded_map,
    entities,
    threshold,
    output_format: str,
) -> dict:
    suffix = Path(filename).suffix.lower()
    stem = Path(filename).stem or "document"

    # output_format == "original" + binary input + redact ON → preserve structure.
    if output_format == "original" and do_redact and suffix in BINARY_ORIGINAL_SUFFIXES:
        return _redact_in_place(suffix, stem, data, loaded_map, entities, threshold)

    # Everything else: extract → optionally redact → render.
    return _extract_and_render(
        suffix, stem, data, do_redact, loaded_map, entities, threshold, output_format,
    )


def _redact_in_place(suffix: str, stem: str, data: bytes, loaded_map, entities, threshold) -> dict:
    """Preserve binary structure (DOCX/PPTX/PDF/EML/MSG) while redacting in place."""
    if suffix == ".docx":
        return _redact_docx_in_place(stem, data, loaded_map, entities, threshold)
    if suffix == ".pptx":
        return _redact_pptx_in_place(stem, data, loaded_map, entities, threshold)
    if suffix == ".pdf":
        return _redact_pdf_to_pdf(stem, data, loaded_map, entities, threshold)
    if suffix in (".eml", ".msg"):
        return _redact_email_in_place(suffix, stem, data, loaded_map, entities, threshold)
    raise HTTPException(status_code=400, detail=f"No in-place redactor for {suffix}.")


def _redact_docx_in_place(stem: str, data: bytes, loaded_map, entities, threshold) -> dict:
    from docx import Document

    r = _build_redactor(loaded_map, entities, threshold)
    with io.BytesIO(data) as buf:
        doc = Document(buf)

    def _redact_paragraph(paragraph) -> None:
        if not paragraph.text:
            return
        new_text = r.redact(paragraph.text).text
        if new_text == paragraph.text:
            return
        if not paragraph.runs:
            paragraph.add_run(new_text)
            return
        paragraph.runs[0].text = new_text
        for run in paragraph.runs[1:]:
            run.text = ""

    for paragraph in doc.paragraphs:
        _redact_paragraph(paragraph)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    _redact_paragraph(paragraph)

    out = io.BytesIO()
    doc.save(out)
    return _binary_response(out.getvalue(), f"{stem}.refined.docx", OUTPUT_MIME["docx"], r.mapping)


def _redact_pptx_in_place(stem: str, data: bytes, loaded_map, entities, threshold) -> dict:
    from pptx import Presentation

    r = _build_redactor(loaded_map, entities, threshold)
    with io.BytesIO(data) as buf:
        prs = Presentation(buf)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text:
                        continue
                    new_text = r.redact(run.text).text
                    if new_text != run.text:
                        run.text = new_text

    out = io.BytesIO()
    prs.save(out)
    return _binary_response(out.getvalue(), f"{stem}.refined.pptx", OUTPUT_MIME["pptx"], r.mapping)


def _redact_pdf_to_pdf(stem: str, data: bytes, loaded_map, entities, threshold) -> dict:
    """PDF → PDF re-typesets via the built-in renderer (layout is not preserved)."""
    import tempfile

    from refinery.formats import _render_output, extract_pdf_text

    with io.BytesIO(data) as buf:
        import pdfplumber
        pages: list[str] = []
        with pdfplumber.open(buf) as pdf:
            for page in pdf.pages:
                pages.append(page.extract_text() or "")
        text = "\n\n".join(pages)

    redacted, mapping = _redact_text(text, loaded_map, entities, threshold)
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tf:
        tmp = Path(tf.name)
    try:
        _render_output(redacted, tmp, "pdf")
        rendered = tmp.read_bytes()
    finally:
        tmp.unlink(missing_ok=True)
    note = ("PDF was extracted to text and re-rendered; the original layout "
            "(fonts, images, columns) is not preserved.")
    return _binary_response(rendered, f"{stem}.refined.pdf", OUTPUT_MIME["pdf"], mapping, note=note)


def _redact_email_in_place(suffix: str, stem: str, data: bytes, loaded_map, entities, threshold) -> dict:
    from email import message_from_bytes
    from email.policy import default as default_policy

    from refinery.email_format import _msg_to_email_message, redact_email_message

    r = _build_redactor(loaded_map, entities, threshold)
    redact_fn = lambda t: r.redact(t).text  # noqa: E731

    note: str | None = None
    if suffix == ".eml":
        msg = message_from_bytes(data, policy=default_policy)
    else:
        import tempfile
        with tempfile.NamedTemporaryFile(suffix=".msg", delete=False) as tf:
            tf.write(data)
            tmp_path = Path(tf.name)
        try:
            msg = _msg_to_email_message(tmp_path)
        finally:
            tmp_path.unlink(missing_ok=True)
        note = "Outlook .msg was converted to a redacted .eml; attachments are not preserved."

    redact_email_message(msg, redact_fn)
    return _binary_response(
        bytes(msg), f"{stem}.refined.eml", OUTPUT_MIME["eml"], r.mapping, note=note,
    )


def _extract_and_render(
    suffix: str,
    stem: str,
    data: bytes,
    do_redact: bool,
    loaded_map,
    entities,
    threshold,
    output_format: str,
) -> dict:
    """Extract text/markdown from any input, optionally redact, render to chosen format."""
    import tempfile

    from refinery.extract import extract_for_output
    from refinery.formats import detect_format

    # If the user picked "original", that means "keep the input format" — but we
    # only get here when the input is text-shaped or redaction is off. Pick a
    # sensible target: same-suffix when we can render it, else markdown.
    target_format = output_format
    if target_format == "original":
        if suffix in (".md", ".markdown"):
            target_format = "md"
        elif suffix in (".txt", ".log", ""):
            target_format = "txt"
        elif suffix == ".html":
            target_format = "html"
        elif suffix == ".pdf":
            target_format = "pdf"
        elif suffix == ".docx":
            target_format = "docx"
        elif suffix == ".pptx":
            target_format = "md"  # PPTX → markdown when reformatting
        else:
            target_format = "txt"

    with tempfile.NamedTemporaryFile(suffix=suffix or ".bin", delete=False) as tf:
        tf.write(data)
        tmp_path = Path(tf.name)
    try:
        fmt = detect_format(tmp_path)
        extracted = extract_for_output(
            tmp_path, fmt, markdown=(target_format != "txt"),
        )
    finally:
        tmp_path.unlink(missing_ok=True)

    if do_redact:
        body, mapping = _redact_text(extracted, loaded_map, entities, threshold)
    else:
        body, mapping = extracted, {}

    return _build_response(
        body, mapping, f"{stem}.refined.{target_format}", target_format,
    )


def _build_response(text_md: str, mapping: dict, filename: str, output_format: str,
                    note: str | None = None) -> dict:
    """Render `text_md` into `output_format` and wrap in JSON."""
    if output_format in ("txt", "md"):
        body = {
            "kind": "text",
            "filename": filename,
            "text": text_md,
            "mapping": mapping,
        }
        if note:
            body["note"] = note
        return body

    import tempfile

    from refinery.formats import _render_output

    suffix = f".{output_format}"
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tf:
        tmp_path = Path(tf.name)
    try:
        _render_output(text_md, tmp_path, output_format)
        rendered = tmp_path.read_bytes()
    finally:
        tmp_path.unlink(missing_ok=True)

    if output_format == "html":
        body = {
            "kind": "text",
            "filename": filename,
            "text": rendered.decode("utf-8", errors="replace"),
            "mapping": mapping,
        }
    else:
        body = {
            "kind": "binary",
            "filename": filename,
            "content_b64": base64.b64encode(rendered).decode("ascii"),
            "content_type": OUTPUT_MIME.get(output_format, "application/octet-stream"),
            "mapping": mapping,
        }
    if note:
        body["note"] = note
    return body


def _binary_response(content: bytes, filename: str, mime: str, mapping: dict,
                     note: str | None = None) -> dict:
    body = {
        "kind": "binary",
        "filename": filename,
        "content_b64": base64.b64encode(content).decode("ascii"),
        "content_type": mime,
        "mapping": mapping,
    }
    if note:
        body["note"] = note
    return body


def run(host: str = "127.0.0.1", port: int = 8000, reload: bool = False) -> None:
    import uvicorn

    if reload:
        uvicorn.run("refinery.web:create_app", host=host, port=port, reload=True, factory=True)
    else:
        uvicorn.run(create_app(), host=host, port=port)
