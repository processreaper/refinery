"""Per-format text/markdown extractors."""

from __future__ import annotations

import html as _html
import re
from email import message_from_bytes
from email.message import EmailMessage
from email.policy import default as default_policy
from pathlib import Path

EMAIL_HEADER_ORDER = ["From", "To", "Cc", "Bcc", "Reply-To", "Subject", "Date"]


def _strip_html(html: str) -> str:
    text = re.sub(r"<\s*br\s*/?\s*>", "\n", html, flags=re.I)
    text = re.sub(r"</\s*p\s*>", "\n\n", text, flags=re.I)
    text = re.sub(r"<[^>]+>", "", text)
    return _html.unescape(text).strip()


def extract_text_passthrough(src: Path) -> str:
    return src.read_text(encoding="utf-8")


def extract_pdf(src: Path) -> str:
    from refinery.formats import extract_pdf_text

    return extract_pdf_text(src)


def extract_pdf_markdown(src: Path) -> str:
    """Extract PDF text with heading guesses based on font size + bold (Markify-style)."""
    import fitz

    doc = fitz.open(str(src))
    parts: list[str] = []
    try:
        for page_num, page in enumerate(doc, 1):
            blocks = page.get_text("dict")["blocks"]
            page_lines: list[str] = []
            for block in blocks:
                if block.get("type") != 0:
                    continue
                for line in block.get("lines", []):
                    spans = line.get("spans", [])
                    if not spans:
                        continue
                    text = "".join(s["text"] for s in spans).strip()
                    if not text:
                        continue
                    max_size = max(s["size"] for s in spans)
                    flags = spans[0].get("flags", 0)
                    bold = bool(flags & 2**4)
                    if max_size >= 16 or (bold and max_size >= 13):
                        page_lines.append(f"## {text}")
                    elif max_size >= 13 or bold:
                        page_lines.append(f"### {text}")
                    else:
                        page_lines.append(text)
            if page_lines:
                parts.append(f"<!-- Page {page_num} -->")
                parts.extend(page_lines)
                parts.append("")
    finally:
        doc.close()
    return "\n".join(parts)


def extract_docx_text(src: Path) -> str:
    from docx import Document

    doc = Document(str(src))
    parts: list[str] = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    if p.text.strip():
                        parts.append(p.text)
    return "\n\n".join(parts)


def extract_docx_markdown(src: Path) -> str:
    from docx import Document

    doc = Document(str(src))
    parts: list[str] = []
    for p in doc.paragraphs:
        text = p.text
        if not text.strip():
            continue
        style = (p.style.name if p.style else "").lower()
        if style == "title":
            parts.append(f"# {text}")
        elif style.startswith("heading 1"):
            parts.append(f"# {text}")
        elif style.startswith("heading 2"):
            parts.append(f"## {text}")
        elif style.startswith("heading 3"):
            parts.append(f"### {text}")
        elif style.startswith("heading 4"):
            parts.append(f"#### {text}")
        elif style.startswith("heading 5"):
            parts.append(f"##### {text}")
        elif style.startswith("heading 6"):
            parts.append(f"###### {text}")
        elif "bullet" in style or "list" in style:
            parts.append(f"- {text}")
        else:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    if p.text.strip():
                        parts.append(p.text)
    return "\n\n".join(parts)


def extract_pptx_markdown(src: Path) -> str:
    """Walk a .pptx into markdown: each slide gets a `## Slide N: title` heading
    followed by indented bullets that mirror the paragraph levels.
    """
    from pptx import Presentation

    prs = Presentation(str(src))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        title_text = ""
        body_parts: list[str] = []

        shapes = sorted(
            slide.shapes, key=lambda s: s.top if s.top is not None else 0,
        )

        for shape in shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue

            is_title = (
                shape.name.lower().startswith("title")
                or (not title_text and shape == shapes[0])
            )

            if is_title and not title_text:
                title_text = text
            else:
                for para in shape.text_frame.paragraphs:
                    para_text = para.text.strip()
                    if not para_text:
                        continue
                    level = para.level or 0
                    indent = "  " * level
                    body_parts.append(f"{indent}- {para_text}")

        lines.append(f"## Slide {i}: {title_text}" if title_text else f"## Slide {i}")
        if body_parts:
            lines.extend(body_parts)
        lines.append("")

    return "\n".join(lines)


def extract_pptx_text(src: Path) -> str:
    """Plain-text fallback: drop the markdown markers from the markdown extract."""
    md = extract_pptx_markdown(src)
    out: list[str] = []
    for line in md.split("\n"):
        if line.startswith("## "):
            out.append(line[3:])
        elif line.lstrip().startswith("- "):
            stripped = line.lstrip()
            indent = line[: len(line) - len(stripped)]
            out.append(f"{indent}{stripped[2:]}")
        else:
            out.append(line)
    return "\n".join(out)


def _email_body(msg: EmailMessage) -> str:
    if msg.is_multipart():
        text_plain: str | None = None
        text_html: str | None = None
        for part in msg.walk():
            if part.is_multipart():
                continue
            ctype = part.get_content_type()
            try:
                content = part.get_content()
            except (LookupError, KeyError, ValueError):
                continue
            if not isinstance(content, str):
                continue
            if ctype == "text/plain" and text_plain is None:
                text_plain = content
            elif ctype == "text/html" and text_html is None:
                text_html = content
        if text_plain:
            return text_plain.strip()
        if text_html:
            return _strip_html(text_html)
        return ""

    ctype = msg.get_content_type()
    try:
        content = msg.get_content()
    except (LookupError, KeyError, ValueError):
        return ""
    if not isinstance(content, str):
        return ""
    if ctype == "text/html":
        return _strip_html(content)
    return content.strip()


def _email_to_string(msg: EmailMessage, markdown: bool) -> str:
    lines: list[str] = []
    for h in EMAIL_HEADER_ORDER:
        v = msg[h]
        if v:
            lines.append(f"**{h}:** {v}" if markdown else f"{h}: {v}")
    if not lines:
        return _email_body(msg)
    if markdown:
        lines.append("")
        lines.append("---")
        lines.append("")
    else:
        lines.append("")
    lines.append(_email_body(msg))
    return "\n".join(lines)


def extract_eml(src: Path, markdown: bool = False) -> str:
    msg = message_from_bytes(src.read_bytes(), policy=default_policy)
    return _email_to_string(msg, markdown)


def extract_msg(src: Path, markdown: bool = False) -> str:
    from refinery.email_format import _msg_to_email_message

    msg = _msg_to_email_message(src)
    return _email_to_string(msg, markdown)


def extract_for_output(src: Path, fmt: str, markdown: bool) -> str:
    """Return a plain-text or markdown rendering of `src` based on its format."""
    if fmt == "pdf":
        return extract_pdf_markdown(src) if markdown else extract_pdf(src)
    if fmt == "docx":
        return extract_docx_markdown(src) if markdown else extract_docx_text(src)
    if fmt == "pptx":
        return extract_pptx_markdown(src) if markdown else extract_pptx_text(src)
    if fmt == "eml":
        return extract_eml(src, markdown=markdown)
    if fmt == "msg":
        return extract_msg(src, markdown=markdown)
    return extract_text_passthrough(src)
