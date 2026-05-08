"""Format adapters: read text from a document, write redacted text back."""

from __future__ import annotations

from pathlib import Path
from typing import Callable

Redact = Callable[[str], str]

OUTPUT_FORMATS = {"txt", "md", "pdf", "docx", "html", "rtf", "odt", "epub"}
PANDOC_ONLY_FORMATS = {"html", "rtf", "odt", "epub"}


def _render_output(text_md: str, dest: Path, output_format: str) -> None:
    """Write `text_md` (markdown text) to `dest` in the requested format."""
    if output_format == "txt":
        dest.write_text(text_md, encoding="utf-8")
        return
    if output_format == "md":
        dest.write_text(text_md, encoding="utf-8")
        return

    from refinery import pandoc_render

    if pandoc_render.can_render(output_format):
        pandoc_render.render(text_md, dest, output_format)
        return

    if output_format in PANDOC_ONLY_FORMATS:
        if not pandoc_render.is_pandoc_available():
            raise RuntimeError(
                f"--output-format {output_format} requires pandoc. "
                "Install it (e.g. `brew install pandoc` on macOS) and retry."
            )
        raise RuntimeError(
            f"pandoc cannot produce {output_format} on this machine"
        )

    if output_format == "pdf":
        from refinery.render import render_markdown_to_pdf

        render_markdown_to_pdf(text_md, dest)
        return
    if output_format == "docx":
        from refinery.render import render_markdown_to_docx

        render_markdown_to_docx(text_md, dest)
        return

    raise RuntimeError(f"unknown output format: {output_format}")


def detect_format(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return "pdf"
    if suffix == ".docx":
        return "docx"
    if suffix in (".pptx", ".ppt"):
        return "pptx"
    if suffix == ".eml":
        return "eml"
    if suffix == ".msg":
        return "msg"
    return "text"


def redact_text_file(src: Path, dest: Path, redact: Redact) -> None:
    text = src.read_text(encoding="utf-8")
    dest.write_text(redact(text), encoding="utf-8")


def extract_pdf_text(src: Path) -> str:
    import pdfplumber

    pages: list[str] = []
    with pdfplumber.open(src) as pdf:
        for page in pdf.pages:
            pages.append(page.extract_text() or "")
    return "\n\n".join(pages)


def redact_pdf_file(src: Path, dest: Path, redact: Redact) -> None:
    text = extract_pdf_text(src)
    _render_output(redact(text), dest, "pdf")


def redact_docx_file(src: Path, dest: Path, redact: Redact) -> None:
    from docx import Document

    doc = Document(str(src))

    def _redact_paragraph(paragraph) -> None:
        if not paragraph.text:
            return
        new_text = redact(paragraph.text)
        if new_text == paragraph.text:
            return
        if not paragraph.runs:
            paragraph.add_run(new_text)
            return
        paragraph.runs[0].text = new_text
        for run in paragraph.runs[1:]:
            run.text = ""

    for paragraph in doc.paragraphs:
        _redact_paragraph(paragraph)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    _redact_paragraph(paragraph)

    doc.save(str(dest))


def redact_pptx_file(src: Path, dest: Path, redact: Redact) -> None:
    """Redact a .pptx in place by walking every text-bearing shape."""
    from pptx import Presentation

    prs = Presentation(str(src))
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text:
                        new_text = redact(run.text)
                        if new_text != run.text:
                            run.text = new_text
    prs.save(str(dest))


def redact_file(
    src: Path,
    dest: Path,
    redact: Redact,
    fmt: str | None = None,
    output_format: str = "original",
) -> None:
    """Redact `src` into `dest`."""
    fmt = fmt or detect_format(src)

    if output_format in OUTPUT_FORMATS:
        from refinery.extract import extract_for_output

        as_markdown = output_format != "txt"
        text = extract_for_output(src, fmt, markdown=as_markdown)
        redacted = redact(text)
        _render_output(redacted, dest, output_format)
        return

    if fmt == "pdf":
        redact_pdf_file(src, dest, redact)
    elif fmt == "docx":
        redact_docx_file(src, dest, redact)
    elif fmt == "pptx":
        redact_pptx_file(src, dest, redact)
    elif fmt == "eml":
        from refinery.email_format import redact_eml_file
        redact_eml_file(src, dest, redact)
    elif fmt == "msg":
        from refinery.email_format import redact_msg_file
        redact_msg_file(src, dest, redact)
    else:
        redact_text_file(src, dest, redact)
